"""
Public API for the AI service.

Use get_quiz_generator() to get the configured provider.
Use extract_text(file_path) to pull text out of any supported document.

Provider is chosen via settings.AI_PROVIDER (default: 'gemini').
Supported: 'gemini' (free), 'claude' (paid).
"""
import os
from django.conf import settings

from .base import (
    BaseQuizGenerator, GenerationRequest, GenerationResult,
    GeneratedQuestion, GeneratedChoice
)


def get_quiz_generator() -> BaseQuizGenerator:
    """Factory: return the configured AI provider. Easy to swap later."""
    provider = (getattr(settings, 'AI_PROVIDER', 'gemini') or 'gemini').lower()

    if provider == 'gemini':
        from .gemini_provider import GeminiQuizGenerator
        return GeminiQuizGenerator()

    if provider == 'claude':
        from .claude_provider import ClaudeQuizGenerator
        return ClaudeQuizGenerator()

    raise ValueError(
        f"Непознат AI провајдер: '{provider}'. "
        f"Поддржани: 'gemini', 'claude'."
    )


def extract_text(file_path: str) -> str:
    """
    Extract plain text from a supported document.
    Returns empty string if the format is unsupported or extraction fails.
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == '.pdf':
            return _extract_pdf(file_path)
        if ext in ('.doc', '.docx'):
            return _extract_docx(file_path)
        if ext in ('.ppt', '.pptx'):
            return _extract_pptx(file_path)
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    except Exception as e:
        print(f"[extract_text] error on {file_path}: {e}")
    return ''


def _extract_pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or '')
        except Exception:
            continue
    return '\n\n'.join(parts).strip()


def _extract_docx(path: str) -> str:
    from docx import Document
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(' | '.join(cells))
    return '\n'.join(parts).strip()


def _extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Слајд {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
    return '\n'.join(parts).strip()


__all__ = [
    'get_quiz_generator',
    'extract_text',
    'BaseQuizGenerator',
    'GenerationRequest',
    'GenerationResult',
    'GeneratedQuestion',
    'GeneratedChoice',
]
